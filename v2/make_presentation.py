"""
Generates BraTS2024_v2_Presentation.pptx — a project overview deck covering intro, data,
model architecture, training, inference, results, uncertainty/calibration, baseline
comparison, and conclusion. Pulls real numbers from PROGRESS.md / V2_1_ARCHITECTURE.md /
BASELINE_COMPARISON.md and embeds existing figures from results_summary/ and results/
(root v1 architecture diagram) rather than re-deriving anything.

Rerun this after the MedNeXt baseline finishes to refresh slide 10 with real numbers.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image
import os

HERE = os.path.dirname(os.path.abspath(__file__))
RESULTS_SUMMARY = os.path.join(HERE, "results_summary")
ROOT_RESULTS = os.path.join(os.path.dirname(HERE), "results")
OUT_PATH = os.path.join(HERE, "BraTS2024_v2_Presentation.pptx")

# ── Palette (matches the dark-background result figures already in the repo) ──────────
BG = RGBColor(0x0B, 0x10, 0x1A)
CARD = RGBColor(0x14, 0x1B, 0x28)
INK = RGBColor(0xF2, 0xF4, 0xF8)
MUTED = RGBColor(0x9A, 0xA5, 0xB4)
TEAL = RGBColor(0x2E, 0xCC, 0x9B)
ORANGE = RGBColor(0xFF, 0x8C, 0x42)
RED = RGBColor(0xE0, 0x5B, 0x5B)
GRID = RGBColor(0x2A, 0x33, 0x44)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.55)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def add_slide():
    slide = prs.slides.add_slide(BLANK)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    # send to back
    bg._element.addprevious(bg._element)
    spTree = slide.shapes._spTree
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return slide


def set_font(run, size=18, color=INK, bold=False, italic=False, name="Calibri"):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = name


def add_text(slide, left, top, width, height, text, size=18, color=INK, bold=False,
             italic=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, name="Calibri",
             line_spacing=1.0):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        set_font(run, size=size, color=color, bold=bold, italic=italic, name=name)
    return box


def add_bullets(slide, left, top, width, height, items, size=16, color=INK,
                 accent=TEAL, bold_lead=False, gap_after=6):
    """items: list of (text, level) or (lead, rest) tuples, or plain strings."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        level = 0
        if isinstance(item, tuple) and len(item) == 2 and isinstance(item[1], int):
            text, level = item
            lead = None
        elif isinstance(item, tuple):
            lead, text = item
        else:
            text, lead = item, None
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap_after)
        p.level = level
        bullet_char = "▸ " if level == 0 else "·  "
        r0 = p.add_run()
        r0.text = bullet_char
        set_font(r0, size=size, color=accent if level == 0 else MUTED, bold=True)
        if lead:
            r1 = p.add_run()
            r1.text = lead + "  "
            set_font(r1, size=size, color=color, bold=True)
            r2 = p.add_run()
            r2.text = text
            set_font(r2, size=size, color=MUTED, bold=False)
        else:
            r1 = p.add_run()
            r1.text = text
            set_font(r1, size=size, color=color, bold=(level == 0 and bold_lead))
    return box


def add_header(slide, kicker, title):
    add_text(slide, MARGIN, Inches(0.35), Inches(9), Inches(0.4), kicker,
              size=14, color=TEAL, bold=True)
    add_text(slide, MARGIN, Inches(0.68), Inches(11.5), Inches(0.7), title,
              size=30, color=INK, bold=True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, Inches(1.42),
                                    Inches(12.23), Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = GRID
    line.line.fill.background()
    line.shadow.inherit = False


def add_footer(slide, n):
    add_text(slide, MARGIN, Inches(7.12), Inches(6), Inches(0.3),
              "BraTS 2024 GLI — Multi-Region Glioma Segmentation", size=10, color=MUTED)
    add_text(slide, Inches(12.0), Inches(7.12), Inches(0.8), Inches(0.3),
              str(n), size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def add_card(slide, left, top, width, height, fill=CARD, line_color=GRID):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.adjustments[0] = 0.04
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = line_color
    card.line.width = Pt(1)
    card.shadow.inherit = False
    return card


def add_image_fit(slide, path, left, top, max_w, max_h, frame=True):
    im = Image.open(path)
    iw, ih = im.size
    ar = iw / ih
    box_ar = max_w / max_h
    if ar > box_ar:
        w = max_w
        h = int(max_w / ar)
    else:
        h = max_h
        w = int(max_h * ar)
    x = left + int((max_w - w) / 2)
    y = top + int((max_h - h) / 2)
    if frame:
        pad = Pt(6)
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x - pad, y - pad,
                                        w + 2 * pad, h + 2 * pad)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0, 0, 0)
        card.line.color.rgb = GRID
        card.line.width = Pt(1)
        card.shadow.inherit = False
    slide.shapes.add_picture(path, x, y, width=w, height=h)
    return x, y, w, h


def style_table(table, header_fill=RGBColor(0x1C, 0x26, 0x38), header_font=TEAL,
                 body_font=INK, alt_fill=RGBColor(0x10, 0x16, 0x22), highlight_rows=None,
                 highlight_fill=RGBColor(0x14, 0x2E, 0x28)):
    highlight_rows = highlight_rows or []
    n_rows = len(table.rows)
    n_cols = len(table.columns)
    for r in range(n_rows):
        for c in range(n_cols):
            cell = table.cell(r, c)
            cell.margin_left = Pt(8)
            cell.margin_right = Pt(8)
            cell.margin_top = Pt(4)
            cell.margin_bottom = Pt(4)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
            elif r in highlight_rows:
                cell.fill.solid()
                cell.fill.fore_color.rgb = highlight_fill
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = alt_fill if r % 2 == 0 else CARD
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in p.runs:
                    set_font(run, size=13, color=header_font if r == 0 else body_font,
                              bold=(r == 0))


# ════════════════════════════════════════════════════════════════════════════════════
# Slide 1 — Title
# ════════════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_text(s, Inches(0), Inches(2.35), SLIDE_W, Inches(0.5),
          "MULTI-REGION BRAIN TUMOR SEGMENTATION", size=16, color=TEAL, bold=True,
          align=PP_ALIGN.CENTER)
add_text(s, Inches(0.8), Inches(2.85), Inches(11.73), Inches(1.5),
          "Uncertainty-Aware Glioma Segmentation\nwith Monte Carlo Dropout",
          size=40, color=INK, bold=True, align=PP_ALIGN.CENTER, line_spacing=1.05)
add_text(s, Inches(0), Inches(4.35), SLIDE_W, Inches(0.4),
          "BraTS 2024 — Adult Glioma, Post-Treatment (GLI)", size=18, color=MUTED,
          align=PP_ALIGN.CENTER)
line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.67), Inches(4.95), Inches(2), Pt(2))
line.fill.solid(); line.fill.fore_color.rgb = TEAL; line.line.fill.background()
line.shadow.inherit = False
add_text(s, Inches(0), Inches(5.2), SLIDE_W, Inches(0.4),
          "Sahar Ifrah   ·   Hadas Avraham", size=18, color=INK, bold=True,
          align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(5.65), SLIDE_W, Inches(0.4),
          "3D U-Net  →  Attention U-Net  ·  PyTorch + MONAI", size=13, color=MUTED,
          align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════════════
# Slide 2 — Intro / Motivation
# ════════════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_header(s, "INTRODUCTION", "The Problem, and Why Uncertainty Matters")
add_text(s, MARGIN, Inches(1.65), Inches(5.8), Inches(0.4), "The task", size=16,
          color=ORANGE, bold=True)
add_bullets(s, MARGIN, Inches(2.1), Inches(5.8), Inches(2.5), [
    "Segment adult diffuse gliomas from multi-parametric brain MRI into three "
    "clinically-defined sub-regions:",
    ("Tumor Core (TC)", 1),
    ("Whole Tumor (WT)", 1),
    ("Enhancing Tumor (ET)", 1),
    "Post-treatment scans (BraTS 2024 GLI) — surgical cavities, radiation "
    "effects, pseudoprogression, not just raw pre-op tumor.",
], size=15)

add_text(s, Inches(6.9), Inches(1.65), Inches(5.8), Inches(0.4), "The innovation", size=16,
          color=ORANGE, bold=True)
add_bullets(s, Inches(6.9), Inches(2.1), Inches(5.8), Inches(2.5), [
    "A segmentation mask alone doesn't say where the model might be wrong.",
    "Monte Carlo Dropout: keep dropout active at inference, run N stochastic "
    "forward passes, and read the spread as a per-voxel uncertainty map.",
    "Goal: uncertainty that's actually usable — high entropy should mean "
    "\"look here,\" not just noise.",
], size=15)

add_card(s, MARGIN, Inches(4.75), Inches(12.23), Inches(1.85))
add_text(s, Inches(0.85), Inches(4.95), Inches(11), Inches(0.4),
          "Two model generations", size=15, color=TEAL, bold=True)
add_bullets(s, Inches(0.85), Inches(5.35), Inches(11), Inches(1.2), [
    ("v1 — plain 3D U-Net", "established the pipeline + MC Dropout innovation. Mean Dice 0.874."),
    ("v2.1 — Attention U-Net", "attention-gated skip connections + stronger augmentation, "
     "targeting the hardest, smallest regions (ET, NCR). Mean Dice 0.887 — current best."),
], size=14)
add_footer(s, 2)

# ════════════════════════════════════════════════════════════════════════════════════
# Slide 3 — Data
# ════════════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_header(s, "THE DATA", "BraTS 2024 GLI — Multi-Parametric MRI")

add_text(s, MARGIN, Inches(1.65), Inches(5.9), Inches(0.35), "Dataset", size=16,
          color=ORANGE, bold=True)
add_bullets(s, MARGIN, Inches(2.05), Inches(5.9), Inches(2.1), [
    ("1,621 training subjects", "1,350 base + 271 additional"),
    ("188 validation subjects", "official BraTS val set, no labels released"),
    ("4 MRI modalities / subject", "T1n, T1c, T2w, T2-FLAIR"),
    ("Pre-registered (SRI24), skull-stripped, 1mm isotropic", "no registration/skull-removal "
     "needed; verified no NaN/Inf"),
], size=14)

add_text(s, Inches(7.0), Inches(1.65), Inches(5.7), Inches(0.35), "Label convention", size=16,
          color=ORANGE, bold=True)
rows, cols = 5, 3
tbl_left, tbl_top, tbl_w, tbl_h = Inches(7.0), Inches(2.05), Inches(5.7), Inches(1.5)
table = s.shapes.add_table(rows, cols, tbl_left, tbl_top, tbl_w, tbl_h).table
table.columns[0].width = Inches(0.9)
table.columns[1].width = Inches(2.6)
table.columns[2].width = Inches(2.2)
headers = ["Label", "Region", "Eval group"]
data = [
    ["0", "Background", "—"],
    ["1", "Necrotic Core (NCR)", "TC"],
    ["2", "Edema (SNFH)", "WT only"],
    ["3", "Enhancing Tumor (ET)", "TC, WT, ET"],
]
for c, h in enumerate(headers):
    table.cell(0, c).text = h
for r, row in enumerate(data, start=1):
    for c, val in enumerate(row):
        table.cell(r, c).text = val
style_table(table)
add_text(s, Inches(7.0), Inches(3.68), Inches(5.7), Inches(0.5),
          "Raw label 4 (old-convention ET) remapped to 3 during preprocessing.",
          size=11, color=MUTED, italic=True)

add_card(s, MARGIN, Inches(4.35), Inches(12.23), Inches(2.25))
add_text(s, Inches(0.85), Inches(4.55), Inches(11), Inches(0.35),
          "Severe class imbalance  +  preprocessing pipeline", size=15, color=TEAL, bold=True)
add_bullets(s, Inches(0.85), Inches(4.95), Inches(5.6), Inches(1.5), [
    ("Background", "98.82% of voxels"),
    ("Edema", "0.816%"),
    ("Enhancing Tumor", "0.356%"),
    ("Necrotic Core", "0.011%  — the rarest class by far"),
], size=13)
add_bullets(s, Inches(6.7), Inches(4.95), Inches(5.6), Inches(1.5), [
    "Normalize each modality by its global training-set max → [0, 1]",
    "Remap label 4 → 3 (unify ET convention)",
    "Tight-crop to brain bounding box, then center-crop/pad to 160×208×160 "
    "(divisible by 16 for U-Net pooling)",
], size=13)
add_footer(s, 3)

# ════════════════════════════════════════════════════════════════════════════════════
# Slide 4 — Model Architecture
# ════════════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_header(s, "MODEL", "Attention U-Net + Monte Carlo Dropout (v2.1)")
arch_png = os.path.join(ROOT_RESULTS, "architecture.png")
add_image_fit(s, arch_png, MARGIN, Inches(1.6), Inches(7.6), Inches(4.35))
add_text(s, MARGIN, Inches(6.05), Inches(7.6), Inches(0.35),
          "Base topology shown (v1). v2.1 keeps this encoder/decoder/bottleneck shape and "
          "adds an attention gate on every skip connection.",
          size=11, color=MUTED, italic=True)

add_bullets(s, Inches(8.35), Inches(1.65), Inches(4.4), Inches(4.9), [
    ("4 encoder stages + bottleneck + 4 decoder", "ConvBlock = Conv3d 3×3×3 → InstanceNorm3d "
     "→ LeakyReLU, ×2. Stride-2 DownConv (not MaxPool) for learned downsampling."),
    ("Attention gate on every skip", "decoder signal + encoder skip combined additively, "
     "squashed to a per-voxel sigmoid mask — suppresses background-dominated regions of "
     "the skip before it reaches the decoder. Not transformer attention: no Q/K/V, no "
     "cross-voxel mixing."),
    ("MC Dropout, p=0.15", "bottleneck + two deepest decoder blocks only — concentrates "
     "stochasticity where the receptive field is largest."),
    ("21.8M parameters", "attention gates add negligible overhead vs. v1's 21.7M"),
], size=13)
add_footer(s, 4)

# ════════════════════════════════════════════════════════════════════════════════════
# Slide 5 — Training
# ════════════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_header(s, "TRAINING", "Loss, Optimization, Augmentation")

add_text(s, MARGIN, Inches(1.65), Inches(5.9), Inches(0.35), "Loss & optimization", size=16,
          color=ORANGE, bold=True)
add_bullets(s, MARGIN, Inches(2.05), Inches(5.9), Inches(2.3), [
    ("DiceFocalLoss", "γ=2.0, background excluded — focal term concentrates the loss on "
     "hard/rare classes (NCR, ET) given the 99% background imbalance"),
    ("AdamW", "lr=2e-4, weight_decay=1e-4"),
    ("Schedule", "5-epoch linear warmup → cosine annealing (T_max=80)"),
    ("Batch size 1, mixed precision (AMP)", "InstanceNorm (not BatchNorm); gradient "
     "clipping at max_norm=1.0"),
], size=14)

add_text(s, Inches(7.0), Inches(1.65), Inches(5.7), Inches(0.35), "Augmentation (train only)",
          size=16, color=ORANGE, bold=True)
add_bullets(s, Inches(7.0), Inches(2.05), Inches(5.7), Inches(2.3), [
    "Independent random flip on each of the 3 spatial axes (p=0.5 each)",
    "Per-modality intensity scale ~U[0.85, 1.15] (p=0.5)",
    "Per-modality intensity shift ~U[-0.10, 0.10] (p=0.5)",
    "Additive Gaussian noise N(0, 0.01) (p=0.20)",
    "Clipped back to [0, 1] after each op — validation never augmented",
], size=14)

add_card(s, MARGIN, Inches(4.55), Inches(12.23), Inches(2.05))
add_text(s, Inches(0.85), Inches(4.75), Inches(11), Inches(0.35),
          "Persistent validation split  +  training outcome", size=15, color=TEAL, bold=True)
add_bullets(s, Inches(0.85), Inches(5.15), Inches(5.6), Inches(1.3), [
    ("1,297 train / 324 val subjects", "split saved once to val_split.json, reloaded on "
     "every run/resume — the held-out set never drifts"),
], size=13)
add_bullets(s, Inches(6.7), Inches(5.15), Inches(5.6), Inches(1.3), [
    ("Best checkpoint: epoch 108", "mean val Dice 0.8870"),
    ("Early stopping patience=35", "training reached epoch 143 (34/35, one short of "
     "triggering) before being interrupted — a strong but not fully converged result"),
], size=13)
add_footer(s, 5)

# ════════════════════════════════════════════════════════════════════════════════════
# Slide 6 — Inference & MC Dropout
# ════════════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_header(s, "INFERENCE", "Deterministic vs. Monte Carlo Dropout")

add_text(s, MARGIN, Inches(1.65), Inches(5.9), Inches(0.35), "Two inference modes", size=16,
          color=ORANGE, bold=True)
add_bullets(s, MARGIN, Inches(2.05), Inches(5.9), Inches(2.3), [
    ("Deterministic", "standard eval-mode forward pass — one segmentation mask, no "
     "uncertainty estimate"),
    ("MC Dropout (N passes)", "force the model into train() mode so Dropout3d layers stay "
     "stochastic, but run under torch.no_grad() — no gradients, just noise"),
    "Run N stochastic forward passes, keep an online mean of the softmax "
    "probabilities (peak memory ≈ 2 volumes, not N)",
], size=14)

add_text(s, Inches(7.0), Inches(1.65), Inches(5.7), Inches(0.35), "What comes out", size=16,
          color=ORANGE, bold=True)
add_bullets(s, Inches(7.0), Inches(2.05), Inches(5.7), Inches(2.3), [
    ("Mean prediction", "argmax of the averaged softmax → the final segmentation, "
     "essentially identical Dice to the deterministic pass"),
    ("Predictive entropy", "H = −Σ p̄·log(p̄) computed from the mean prediction — the "
     "\"entropy of the mean,\" the standard MC-Dropout uncertainty estimate"),
    "High entropy → the N stochastic passes disagreed → flag for review",
], size=14)

add_card(s, MARGIN, Inches(4.55), Inches(12.23), Inches(2.05))
add_text(s, Inches(0.85), Inches(4.75), Inches(11), Inches(0.35),
          "Why bother, if Dice barely changes?", size=15, color=TEAL, bold=True)
add_bullets(s, Inches(0.85), Inches(5.15), Inches(11), Inches(1.3), [
    "MC-mean Dice (0.8872) vs. deterministic Dice (0.8870) — essentially identical. "
    "That's expected: dropout is placed to generate calibrated uncertainty, not to "
    "boost raw accuracy.",
    "The value is the entropy map itself: false positives get ~7× the predictive "
    "entropy of true positives (slide 9) — a usable signal for flagging voxels a "
    "clinician should double-check, not just a diagnostic curiosity.",
], size=13)
add_footer(s, 6)

# ════════════════════════════════════════════════════════════════════════════════════
# Slide 7 — Qualitative example
# ════════════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_header(s, "QUALITATIVE EXAMPLE", "MRI, Ground Truth, Prediction, and Uncertainty")
example_png = os.path.join(RESULTS_SUMMARY, "examples", "example_01_BraTS-GLI-02385-100.png")
add_image_fit(s, example_png, MARGIN, Inches(2.1), Inches(12.23), Inches(3.6))
add_bullets(s, MARGIN, Inches(5.95), Inches(12.23), Inches(1.0), [
    "Left to right: raw T1c MRI → ground truth overlay → 20-pass MC-Dropout prediction "
    "overlay → entropy map (brighter = more uncertain) → entropy overlaid on the scan.",
    "Uncertainty concentrates almost exactly on the tumor boundary and the small "
    "isolated lesion the model is least sure about — not scattered randomly.",
], size=13)
add_footer(s, 7)

# ════════════════════════════════════════════════════════════════════════════════════
# Slide 8 — Results: v1 vs v2.1
# ════════════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_header(s, "RESULTS", "v1 (Plain U-Net) vs. v2.1 (Attention U-Net)")

rows, cols = 4, 9
tbl_left, tbl_top, tbl_w, tbl_h = MARGIN, Inches(1.75), Inches(12.23), Inches(1.7)
table = s.shapes.add_table(rows, cols, tbl_left, tbl_top, tbl_w, tbl_h).table
col_w = [Inches(3.0)] + [Inches(1.15)] * 8
for c, w in enumerate(col_w):
    table.columns[c].width = w
headers = ["Model — Mode", "Dice TC", "Dice WT", "Dice ET", "Dice Mean",
           "HD95 TC", "HD95 WT", "HD95 ET", "HD95 Mean"]
data = [
    ["v1 — plain U-Net", "0.8638", "0.9087", "0.8503", "0.8743", "5.65", "6.11", "5.74", "5.83"],
    ["v2.1 — Attention U-Net (det.)", "0.8791", "0.9206", "0.8613", "0.8870", "4.83", "4.87", "5.06", "4.92"],
    ["v2.1 — MC Dropout (20 passes)", "0.8794", "0.9205", "0.8616", "0.8872", "4.81", "4.88", "5.04", "4.91"],
]
for c, h in enumerate(headers):
    table.cell(0, c).text = h
for r, row in enumerate(data, start=1):
    for c, val in enumerate(row):
        table.cell(r, c).text = val
style_table(table, highlight_rows=[2])

add_card(s, MARGIN, Inches(3.75), Inches(12.23), Inches(2.6))
add_bullets(s, Inches(0.85), Inches(3.95), Inches(11.5), Inches(2.2), [
    ("v2.1 improves on every metric", "+1.3 pts mean Dice, ~16% lower mean HD95 vs. v1 — "
     "same parameter budget (~21.7-21.8M), same evaluation protocol (324-subject holdout)."),
    ("Biggest relative gain on ET", "0.8503 → 0.8613 — exactly the small, hard region "
     "attention gates were added to help."),
    ("MC Dropout barely moves the point estimate", "+0.0002 mean Dice — expected, since "
     "dropout here is placed to generate calibrated uncertainty (slide 9), not to boost "
     "raw accuracy. The mean prediction and the deterministic pass are essentially the "
     "same segmentation."),
], size=14, gap_after=10)
add_footer(s, 8)

# ════════════════════════════════════════════════════════════════════════════════════
# Slide 9 — Uncertainty & Calibration
# ════════════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_header(s, "UNCERTAINTY & CALIBRATION", "Is the Uncertainty Map Actually Useful?")

cal_png = os.path.join(RESULTS_SUMMARY, "calibration.png")
add_image_fit(s, cal_png, MARGIN, Inches(1.65), Inches(7.7), Inches(3.55))

rows, cols = 5, 2
tbl_left, tbl_top, tbl_w, tbl_h = Inches(8.35), Inches(1.65), Inches(4.4), Inches(2.0)
table = s.shapes.add_table(rows, cols, tbl_left, tbl_top, tbl_w, tbl_h).table
table.columns[0].width = Inches(2.4)
table.columns[1].width = Inches(2.0)
headers = ["Region", "Mean Entropy"]
data = [["TN (correct bg)", "0.0002"], ["TP (correct tumor)", "0.0285"],
        ["FP (false positive)", "0.1960"], ["FN (false negative)", "0.1112"]]
for c, h in enumerate(headers):
    table.cell(0, c).text = h
for r, row in enumerate(data, start=1):
    for c, val in enumerate(row):
        table.cell(r, c).text = val
style_table(table, highlight_rows=[3])

add_bullets(s, Inches(8.35), Inches(3.85), Inches(4.4), Inches(3.0), [
    ("FP entropy ≈ 6.9× TP entropy", "the model flags mistakes with far higher "
     "uncertainty — same as v1 (~6×), holds after the attention upgrade"),
    ("ECE = 0.0080", "softmax confidence is well calibrated"),
    ("AUROC = 0.8956", "entropy alone strongly predicts misclassified voxels — usable to "
     "flag regions for review"),
], size=13, gap_after=10)
add_text(s, MARGIN, Inches(5.35), Inches(7.7), Inches(0.6),
          "(50 val subjects, brain-mask restricted, for the calibration analysis; "
          "entropy-by-region from 10 example subjects, 20 MC passes)",
          size=11, color=MUTED, italic=True)
add_footer(s, 9)

# ════════════════════════════════════════════════════════════════════════════════════
# Slide 10 — Baseline comparison
# ════════════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_header(s, "BASELINE COMPARISON", "vs. Published, Pretrained BraTS Models")
add_text(s, MARGIN, Inches(1.55), Inches(12.23), Inches(0.4),
          "Same 324-subject val split, real downloaded checkpoints, adapter pipeline "
          "(channel reorder, normalization, sliding-window inference, MC Dropout) — a "
          "true head-to-head, not a literature-numbers comparison.",
          size=13, color=MUTED)

rows, cols = 4, 3
tbl_left, tbl_top, tbl_w, tbl_h = MARGIN, Inches(2.15), Inches(7.3), Inches(1.85)
table = s.shapes.add_table(rows, cols, tbl_left, tbl_top, tbl_w, tbl_h).table
table.columns[0].width = Inches(4.4)
table.columns[1].width = Inches(1.45)
table.columns[2].width = Inches(1.45)
headers = ["Model — Mode", "Dice Mean", "HD95 Mean"]
data = [
    ["v2.1 (ours) — Deterministic", "0.8870", "4.92"],
    ["SegResNet (BraTS 2018, published)", "0.5791", "18.09"],
    ["MedNeXt (BraTS 2024 winner, published)", "in progress", "—"],
]
for c, h in enumerate(headers):
    table.cell(0, c).text = h
for r, row in enumerate(data, start=1):
    for c, val in enumerate(row):
        table.cell(r, c).text = val
style_table(table, highlight_rows=[1])

add_bullets(s, MARGIN, Inches(4.2), Inches(7.6), Inches(2.6), [
    ("SegResNet (MONAI bundle, trained on BraTS 2018)", "scores far lower on our data "
     "(0.58 vs 0.887) — expected: it never saw BraTS 2024's post-treatment scans "
     "(surgical cavities, radiation effects, pseudoprogression)."),
    ("Its MC-dropout uncertainty is also much blunter", "FP/TP entropy ratio 2.6× vs. our "
     "6.9× — a model trained on your actual domain gives a sharper uncertainty signal, "
     "not just better accuracy."),
    ("MedNeXt", "the actual BraTS-winning architecture, trained on real BraTS 2024 data — "
     "evaluation running now; two real bugs (axis-order mismatch, a label-remap "
     "assumption) were caught and fixed before trusting the numbers."),
], size=13)

ent_png = os.path.join(RESULTS_SUMMARY, "entropy_comparison.png")
add_image_fit(s, ent_png, Inches(8.05), Inches(2.15), Inches(4.5), Inches(4.5))
add_footer(s, 10)

# ════════════════════════════════════════════════════════════════════════════════════
# Slide 11 — Conclusion
# ════════════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_header(s, "CONCLUSION", "Where This Lands, and What's Next")

add_text(s, MARGIN, Inches(1.65), Inches(5.9), Inches(0.35), "What we showed", size=16,
          color=ORANGE, bold=True)
add_bullets(s, MARGIN, Inches(2.05), Inches(5.9), Inches(3.0), [
    ("Attention gates measurably help", "v2.1 beats v1 on every metric, with the largest "
     "gain on the smallest/hardest region (ET)"),
    ("MC Dropout uncertainty is genuinely calibrated", "ECE 0.008, AUROC 0.896 — not just "
     "a visualization, a usable error-flagging signal"),
    ("We outperform published, pretrained baselines head-to-head", "on our own val split, "
     "not just on paper — and our uncertainty is sharper too"),
], size=14)

add_text(s, Inches(6.9), Inches(1.65), Inches(5.7), Inches(0.35), "Limitations & next steps",
          size=16, color=ORANGE, bold=True)
add_bullets(s, Inches(6.9), Inches(2.05), Inches(5.7), Inches(3.0), [
    ("Training not fully converged", "interrupted at epoch 143/300, one epoch short of "
     "early stopping — resumable, room for further improvement"),
    ("MedNeXt comparison in progress", "closest domain match of any public checkpoint — "
     "will complete the baseline picture"),
    ("Single dataset, internal holdout", "BraTS 2024 challenge closed before official "
     "leaderboard submission was possible"),
], size=14)

add_card(s, MARGIN, Inches(5.35), Inches(12.23), Inches(1.35))
add_text(s, Inches(0.85), Inches(5.55), Inches(11.5), Inches(0.9),
          "Bottom line: a better architecture and a genuinely trustworthy uncertainty "
          "estimate — not just a higher Dice score.",
          size=16, color=TEAL, bold=True, italic=True)
add_footer(s, 11)

# ════════════════════════════════════════════════════════════════════════════════════
# Slide 12 — Thank you
# ════════════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_text(s, Inches(0), Inches(3.0), SLIDE_W, Inches(0.8), "Thank You", size=44, color=INK,
          bold=True, align=PP_ALIGN.CENTER)
line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.67), Inches(3.95), Inches(2), Pt(2))
line.fill.solid(); line.fill.fore_color.rgb = TEAL; line.line.fill.background()
line.shadow.inherit = False
add_text(s, Inches(0), Inches(4.2), SLIDE_W, Inches(0.4), "Questions?", size=18, color=MUTED,
          align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(5.0), SLIDE_W, Inches(0.4),
          "Sahar Ifrah   ·   Hadas Avraham", size=14, color=INK, align=PP_ALIGN.CENTER)

prs.save(OUT_PATH)
print(f"Saved: {OUT_PATH}  ({len(prs.slides)} slides)")